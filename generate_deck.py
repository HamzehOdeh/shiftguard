"""
ShiftGuard Investor/Sales Pitch Deck Generator
Generates a professional PowerPoint with branded visuals.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BRAND_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
BRAND_DARK = RGBColor(0x03, 0x07, 0x12)
BRAND_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
RED_400 = RGBColor(0xF8, 0x71, 0x71)
GREEN_400 = RGBColor(0x4A, 0xDE, 0x80)
ORANGE_400 = RGBColor(0xFB, 0x92, 0x3C)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, text_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color or color
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_stat_card(slide, left, top, number, label, color=BRAND_BLUE):
    width = Inches(2.2)
    height = Inches(1.4)
    shape = add_rounded_rect(slide, left, top, width, height, BRAND_NAVY)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_400
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== SLIDE 1: Title =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    # Logo square
    logo = add_rounded_rect(slide, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5), BRAND_BLUE, "SG", font_size=36)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2),
                 "ShiftGuard", font_size=48, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.6), Inches(9.3), Inches(0.8),
                 "AI-Powered Workforce Compliance Intelligence", font_size=24, color=GRAY_300, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.6),
                 "Schedule with confidence. Never pay another compliance fine.",
                 font_size=16, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(4), Inches(6.5), Inches(5.3), Inches(0.5),
                 "July 2026 | Confidential", font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    # ===================== SLIDE 2: The Problem =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "The Problem", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
                 "U.S. employers pay $2.7 BILLION in wage & hour penalties every year.", font_size=20, color=GRAY_300)

    # Stat cards
    add_stat_card(slide, Inches(0.8), Inches(2.2), "$2.7B", "Annual penalties\n(DOL enforcement)", RED_400)
    add_stat_card(slide, Inches(3.3), Inches(2.2), "70%", "DOL investigations\nfind violations", ORANGE_400)
    add_stat_card(slide, Inches(5.8), Inches(2.2), "12+", "New scheduling laws\nsince 2022", BRAND_BLUE)
    add_stat_card(slide, Inches(8.3), Inches(2.2), "4-8 hrs", "Manager time/week\nchecking compliance", GRAY_300)
    add_stat_card(slide, Inches(10.8), Inches(2.2), "$50M+", "Class-action\nsettlements", RED_400)

    # Problem bullets
    lines = [
        ("The root cause isn't malice — it's complexity:", True, WHITE),
        ("", False, None),
        ("  A single warehouse in Chicago must comply with FLSA + Illinois ODRISA +", False, GRAY_300),
        ("  Chicago Fair Workweek + union CBA + company policy + OSHA — simultaneously,", False, GRAY_300),
        ("  for every shift, for every worker.", False, GRAY_300),
        ("", False, None),
        ("  Existing tools catch violations AFTER they happen.", False, GRAY_300),
        ("  By then, the penalty is already incurred.", False, GRAY_300),
        ("", False, None),
        ("  Managers don't know the laws. HR updates lag. Nobody checks in real-time.", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(3.2), lines, font_size=15)

    # ===================== SLIDE 3: The Solution =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "The Solution", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
                 "Check every schedule against every law BEFORE you publish.", font_size=20, color=BRAND_BLUE)

    # Flow diagram
    steps = [
        ("Schedule\nCreated", BRAND_NAVY),
        ("ShiftGuard\nScans", BRAND_BLUE),
        ("Violations\nFlagged ($)", RGBColor(0x7C, 0x3A, 0xED)),
        ("AI Fix\nSuggested", RGBColor(0x06, 0x94, 0xA2)),
        ("Publish with\nConfidence", GREEN_400),
    ]
    for i, (text, color) in enumerate(steps):
        left = Inches(0.8 + i * 2.5)
        add_rounded_rect(slide, left, Inches(2.3), Inches(2.1), Inches(1.2), color, text, font_size=13)
        if i < len(steps) - 1:
            add_text_box(slide, left + Inches(2.1), Inches(2.7), Inches(0.4), Inches(0.4),
                         "→", font_size=24, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5), Inches(3.7), Inches(3.3), Inches(0.4),
                 "60 seconds. Zero configuration.", font_size=13, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    # Key capabilities
    lines = [
        ("What ShiftGuard does that nobody else can:", True, WHITE),
        ("", False, None),
        ("  ✔  Shows penalty $ exposure on every violation before publish", False, GREEN_400),
        ("  ✔  AI suggests the least-disruptive compliant fix (one click)", False, GREEN_400),
        ("  ✔  Plain English explanation of every decision", False, GREEN_400),
        ("  ✔  Fairness engine ensures equitable shift distribution", False, GREEN_400),
        ("  ✔  Deploys in 60 seconds — not 6 months", False, GREEN_400),
        ("  ✔  Updates rules within 48 hours of new legislation", False, GREEN_400),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(3), lines, font_size=15)

    # ===================== SLIDE 4: Product Demo =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Product — Live Today", font_size=32, bold=True)

    # Manager side
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "Manager Dashboard", font_size=16, bold=True, color=BRAND_BLUE)

    mgr_lines = [
        ("⚠  CRITICAL: Dr. Patel exceeded 80hr/week (ACGME-001)", False, RED_400),
        ("     Penalty: $7,500/week  |  Fix: Remove Friday night shift", False, GRAY_400),
        ("", False, None),
        ("⚠  HIGH: RN Chen — clopening, <8h rest (REST-002)", False, ORANGE_400),
        ("     Penalty: $5,000  |  Fix: Move morning shift to 10AM", False, GRAY_400),
        ("", False, None),
        ("⚠  MEDIUM: Schedule posted 5 days early (NOTICE-001)", False, RGBColor(0xFB, 0xBF, 0x24)),
        ("     Penalty: $500/worker  |  Fix: Publish by Wednesday", False, GRAY_400),
        ("", False, None),
        ("  Confidence Score: 72%  |  Total Exposure: $18,500/week", False, WHITE),
        ("  [Resolve All]  [Export Report]  [Publish Anyway]", False, BRAND_BLUE),
    ]
    add_multiline_text(slide, Inches(1.1), Inches(2.1), Inches(5.3), Inches(4.5), mgr_lines, font_size=12)

    # Worker side
    add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
                 "Worker Mobile App (PWA)", font_size=16, bold=True, color=BRAND_BLUE)

    worker_lines = [
        ("Good morning, Sarah", True, WHITE),
        ("ED Nursing | Staff RN", False, GRAY_500),
        ("", False, None),
        ("\U0001f319  Night shift in 24h — starts at 19:00", False, ORANGE_400),
        ("    Tip: Nap 90min between 2-4 PM for peak alertness", False, GRAY_400),
        ("", False, None),
        ("⏱  Only 4h before overtime (36/40h this week)", False, BRAND_BLUE),
        ("", False, None),
        ("\U0001f4b5  VET Available: Thu Jul 10, premium pay", False, GREEN_400),
        ("    [Accept]  [Decline]", False, GRAY_400),
        ("", False, None),
        ("✅  PTO Jul 15-17 auto-approved. Enjoy!", False, GREEN_400),
        ("", False, None),
        ("  \U0001f4c5 Schedule  |  \U0001f4dd Request  |  \U0001f4b0 Balance  |  \U0001f4ac AI", False, GRAY_500),
    ]
    add_multiline_text(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(4.5), worker_lines, font_size=11)

    # ===================== SLIDE 5: Why Different =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Why ShiftGuard Wins", font_size=32, bold=True)

    # Comparison table
    headers = ["Capability", "ShiftGuard", "UKG", "Legion", "Deputy"]
    rows = [
        ["Pre-publish compliance check", "✔", "✖", "✖", "✖"],
        ["Penalty $ on every violation", "✔", "✖", "✖", "✖"],
        ["AI fix suggestions", "✔", "✖", "Partial", "✖"],
        ["Multi-jurisdiction auto-detect", "✔", "Manual", "✖", "✖"],
        ["Union CBA enforcement", "✔", "Partial", "✖", "✖"],
        ["Fairness engine", "✔", "✖", "Cost only", "✖"],
        ["Plain English decisions", "✔", "✖", "✖", "✖"],
        ["Setup time", "60 sec", "6-12 mo", "3+ mo", "2-4 wk"],
        ["Price (PEPM)", "$15-25", "$50-100+", "$25-50+", "$6.50"],
    ]

    col_widths = [Inches(3.5), Inches(2), Inches(2), Inches(2), Inches(2)]
    start_top = Inches(1.5)
    row_height = Inches(0.5)
    start_left = Inches(0.8)

    # Header row
    x = start_left
    for i, h in enumerate(headers):
        color = BRAND_BLUE if i == 1 else GRAY_400
        bold = i <= 1
        add_text_box(slide, x, start_top, col_widths[i], Inches(0.5), h, font_size=12, bold=bold, color=color)
        x += col_widths[i]

    # Data rows
    for r_idx, row in enumerate(rows):
        y = start_top + Inches(0.5) + row_height * r_idx
        x = start_left
        for c_idx, cell in enumerate(row):
            if cell == "✔":
                color = GREEN_400
            elif cell == "✖":
                color = RED_400
            elif c_idx == 1:
                color = GREEN_400
            else:
                color = GRAY_300
            add_text_box(slide, x, y, col_widths[c_idx], row_height, cell, font_size=12, color=color)
            x += col_widths[c_idx]

    # ===================== SLIDE 6: Market Opportunity =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Market Opportunity", font_size=32, bold=True)

    # TAM/SAM/SOM circles
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(4.5), Inches(0.5),
                 "TAM: $2.7B  |  SAM: $1.08B  |  SOM (Yr1): $570K", font_size=16, color=BRAND_BLUE)

    # Market segments
    segments = [
        ("Healthcare", "5.3M workers", "$950M", "Hospitals, clinics, health systems"),
        ("Warehousing", "1.8M workers", "$540M", "Fulfillment, logistics, 3PL"),
        ("Retail", "4.2M workers", "$630M", "Multi-location chains"),
        ("Hospitality", "3.1M workers", "$465M", "Hotels, restaurants, venues"),
        ("Manufacturing", "1.2M workers", "$180M", "Plants, factories"),
    ]

    for i, (name, workers, revenue, desc) in enumerate(segments):
        top = Inches(2.2 + i * 1.0)
        add_rounded_rect(slide, Inches(0.8), top, Inches(6), Inches(0.8), BRAND_NAVY)
        add_text_box(slide, Inches(1.0), top + Inches(0.05), Inches(2), Inches(0.4), name, font_size=14, bold=True)
        add_text_box(slide, Inches(1.0), top + Inches(0.4), Inches(4), Inches(0.4), f"{workers} | {desc}", font_size=11, color=GRAY_400)
        add_text_box(slide, Inches(5.2), top + Inches(0.15), Inches(1.5), Inches(0.5), revenue, font_size=16, bold=True, color=GREEN_400, alignment=PP_ALIGN.RIGHT)

    # Why now
    add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5), Inches(0.5),
                 "Why Now?", font_size=20, bold=True, color=BRAND_BLUE)

    why_now = [
        ("12+ new predictive scheduling laws (2022-2026)", False, GRAY_300),
        ("DOL hired 600+ new investigators (2024-25)", False, GRAY_300),
        ("AI maturity: LLMs can reason about law", False, GRAY_300),
        ("Gen Z demands schedule fairness/transparency", False, GRAY_300),
        ("Labor scarcity: retention > recruitment", False, GRAY_300),
        ("$5-50M class actions now routine", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(7.5), Inches(2.0), Inches(5), Inches(4), why_now, font_size=14, spacing=1.8)

    # ===================== SLIDE 7: Business Model =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Business Model", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.5),
                 "SaaS: $15-25 Per Employee Per Month", font_size=20, color=BRAND_BLUE)

    # Tiers
    tiers = [
        ("Core\n$15 PEPM", "Pre-publish scan\nViolation alerts\nFix suggestions\nAudit trail", BRAND_NAVY),
        ("Growth\n$20 PEPM", "+ Fairness engine\n+ Worker portal\n+ Auto-approval\n+ Coverage AI", RGBColor(0x1E, 0x3A, 0x5F)),
        ("Enterprise\n$25 PEPM", "+ AI assistant\n+ Predictive callouts\n+ Demand sensing\n+ Multi-site + SSO", RGBColor(0x0C, 0x4A, 0x6E)),
    ]

    for i, (title, features, color) in enumerate(tiers):
        left = Inches(0.8 + i * 4.1)
        shape = add_rounded_rect(slide, left, Inches(2.0), Inches(3.8), Inches(2.2), color)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BRAND_BLUE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        for line in features.split("\n"):
            p2 = tf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(12)
            p2.font.color.rgb = GRAY_300
            p2.font.name = "Segoe UI"
            p2.alignment = PP_ALIGN.CENTER

    # Revenue projections
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.5),
                 "Revenue Projections", font_size=18, bold=True)

    proj_data = [
        ("Year 1", "15 customers", "$570K ARR", "85% margin"),
        ("Year 2", "50 customers", "$2.5M ARR", "87% margin"),
        ("Year 3", "150 customers", "$8.3M ARR", "90% margin"),
    ]

    for i, (year, custs, arr, margin) in enumerate(proj_data):
        left = Inches(0.8 + i * 4.1)
        top = Inches(5.1)
        add_rounded_rect(slide, left, top, Inches(3.8), Inches(1.6), BRAND_NAVY)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(3.4), Inches(0.4), year, font_size=13, color=GRAY_400, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.45), Inches(3.4), Inches(0.5), arr, font_size=22, bold=True, color=GREEN_400, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), top + Inches(1.05), Inches(3.4), Inches(0.4), f"{custs} | {margin}", font_size=12, color=GRAY_400, alignment=PP_ALIGN.CENTER)

    # ===================== SLIDE 8: GTM Strategy =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Go-To-Market", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
                 "Product-Led Growth: Free tool → Demo → Trial → Close", font_size=18, color=BRAND_BLUE)

    # Funnel
    funnel_steps = [
        ("Free Calculator", "No signup. Instant risk grade.\nShareable. Every HR leader\nwho sees $47K risk = warm lead.", "5,000 visitors/mo", RGBColor(0x06, 0x94, 0xA2)),
        ("Interactive Demo", "Pick industry, see real violations\nwith penalty $. One-click CTA.", "500 demos/mo", BRAND_BLUE),
        ("14-Day Trial", "Upload YOUR schedule.\nSee YOUR violations.\nProve ROI before buying.", "50 trials/mo", RGBColor(0x7C, 0x3A, 0xED)),
        ("Close", "$15-25 PEPM\nAvg deal: $5K/mo\nPayback: <2 months", "5 closes/mo", GREEN_400),
    ]

    for i, (title, desc, metric, color) in enumerate(funnel_steps):
        top = Inches(2.0 + i * 1.3)
        width = Inches(11.5 - i * 1.5)
        add_rounded_rect(slide, Inches(0.8), top, width, Inches(1.1), BRAND_NAVY)
        add_text_box(slide, Inches(1.1), top + Inches(0.1), Inches(2.5), Inches(0.4), title, font_size=14, bold=True, color=color)
        add_text_box(slide, Inches(3.8), top + Inches(0.15), Inches(5), Inches(0.9), desc, font_size=11, color=GRAY_300)
        add_text_box(slide, width - Inches(1.5), top + Inches(0.3), Inches(2.2), Inches(0.4), metric, font_size=12, color=GRAY_500, alignment=PP_ALIGN.RIGHT)

    # Vertical focus
    add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
                 "Phase 1: Healthcare (highest pain)  →  Phase 2: Warehouse  →  Phase 3: Multi-location retail",
                 font_size=14, color=GRAY_400)

    # ===================== SLIDE 9: Moat =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Competitive Moat", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
                 "Five layers of defensibility — each harder to replicate than the last", font_size=16, color=GRAY_300)

    moats = [
        ("Compliance-as-Code", "Every law = versioned executable rule. 50 states + local. Updated in 48h.", "12-18 months to replicate"),
        ("Data Flywheel", "Every customer improves predictions for all. Callouts, demand, violations.", "Requires our scale first"),
        ("Fairness Memory", "Year-over-year allocation history per worker. Switching = losing years of data.", "Cannot replicate without history"),
        ("AI Reasoning", "LLM fine-tuned on scheduling domain. Explains, generates, predicts.", "Commodity AI without domain data"),
        ("Multi-Jurisdiction Speed", "New law → rule deployed in 48h. Competitors: 3-6 months.", "Architecture built for speed"),
    ]

    for i, (title, desc, replicate) in enumerate(moats):
        top = Inches(1.9 + i * 1.1)
        add_rounded_rect(slide, Inches(0.8), top, Inches(0.3), Inches(0.9), BRAND_BLUE)
        add_text_box(slide, Inches(1.3), top + Inches(0.0), Inches(5), Inches(0.4), title, font_size=15, bold=True)
        add_text_box(slide, Inches(1.3), top + Inches(0.4), Inches(6), Inches(0.4), desc, font_size=12, color=GRAY_300)
        add_text_box(slide, Inches(8.0), top + Inches(0.2), Inches(4.5), Inches(0.4), replicate, font_size=12, color=GRAY_500)

    # ===================== SLIDE 10: Traction =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Traction — Live Product", font_size=32, bold=True)

    add_stat_card(slide, Inches(0.8), Inches(1.5), "24+", "Modules built")
    add_stat_card(slide, Inches(3.3), Inches(1.5), "31", "Compliance rules")
    add_stat_card(slide, Inches(5.8), Inches(1.5), "5", "Industries covered")
    add_stat_card(slide, Inches(8.3), Inches(1.5), "50", "US states")
    add_stat_card(slide, Inches(10.8), Inches(1.5), "3", "Live deployments")

    # What's built
    built_lines = [
        ("What's Live:", True, BRAND_BLUE),
        ("", False, None),
        ("  ✔  Full compliance engine (31 rules, 7 violation types, penalty $ per state)", False, GREEN_400),
        ("  ✔  Healthcare app (ACGME, residency scheduler, jeopardy, nursing ratios)", False, GREEN_400),
        ("  ✔  AI assistant \"Otto\" (Claude-powered, answers in plain English)", False, GREEN_400),
        ("  ✔  Worker mobile PWA (installable, calendar, requests, balances)", False, GREEN_400),
        ("  ✔  Free compliance calculator (viral growth tool)", False, GREEN_400),
        ("  ✔  Manager dashboard (violations, ROI, fairness scores)", False, GREEN_400),
        ("  ✔  Data connectors (UKG/Kronos API, Google Sheets, database)", False, GREEN_400),
        ("  ✔  Deployed: Streamlit Cloud + Vercel + GitHub", False, GREEN_400),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(4), built_lines, font_size=14)

    # URLs
    url_lines = [
        ("Live URLs:", True, BRAND_BLUE),
        ("", False, None),
        ("Healthcare Demo:  shiftguard-giznml7xnt59j5aguqjerc.streamlit.app", False, GRAY_300),
        ("Worker App:  frontend-lyart-three-52.vercel.app/worker", False, GRAY_300),
        ("Calculator:  frontend-lyart-three-52.vercel.app/calculator", False, GRAY_300),
        ("Demo:  frontend-lyart-three-52.vercel.app/demo", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(8.5), Inches(3.2), Inches(4.5), Inches(3.5), url_lines, font_size=11)

    # ===================== SLIDE 11: Team =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "The Team", font_size=32, bold=True)

    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.5), BRAND_NAVY)

    add_text_box(slide, Inches(1.5), Inches(1.7), Inches(8), Inches(0.6),
                 "Hamzeh Odeh — Founder & CEO", font_size=22, bold=True)

    founder_lines = [
        ("5 years Amazon Operations Finance — managed compliance and labor cost", False, GRAY_300),
        ("optimization across millions of labor hours at the world's largest employer.", False, GRAY_300),
        ("", False, None),
        ("Why this founder wins:", True, WHITE),
        ("", False, None),
        ("  •  Lived the problem at Amazon scale (1.5M hourly workers)", False, GRAY_300),
        ("  •  Understands both ops pain (manager) and financial impact (finance)", False, GRAY_300),
        ("  •  Technical enough to build, business enough to sell", False, GRAY_300),
        ("  •  Network: ops leaders at companies most likely to buy", False, GRAY_300),
        ("  •  Built the entire platform solo in 6 weeks", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(2.5), founder_lines, font_size=14)

    # Hiring plan
    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(6), Inches(0.5),
                 "First Hires (Post-Funding):", font_size=16, bold=True, color=BRAND_BLUE)

    hire_lines = [
        ("  1. Senior Full-Stack Engineer (Next.js + Python)", False, GRAY_300),
        ("  2. Healthcare Domain Expert / Sales (hospital relationships)", False, GRAY_300),
        ("  3. Compliance Analyst (law library expansion + updates)", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(5.8), Inches(8), Inches(1.5), hire_lines, font_size=14)

    # ===================== SLIDE 12: The Ask =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "The Ask", font_size=32, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6),
                 "Raising $1.5M Seed", font_size=24, bold=True, color=BRAND_BLUE)

    # Use of funds
    funds = [
        ("Engineering (2 senior hires)", "45%", Inches(5.0)),
        ("Sales & GTM (1 AE + marketing)", "30%", Inches(3.3)),
        ("Legal (rule library expansion)", "15%", Inches(1.7)),
        ("Infrastructure & ops", "10%", Inches(1.1)),
    ]

    for i, (label, pct, bar_width) in enumerate(funds):
        top = Inches(2.3 + i * 0.9)
        add_text_box(slide, Inches(0.8), top, Inches(4.5), Inches(0.4), label, font_size=14, color=GRAY_300)
        add_rounded_rect(slide, Inches(5.5), top + Inches(0.05), bar_width, Inches(0.35), BRAND_BLUE)
        add_text_box(slide, Inches(5.5) + bar_width + Inches(0.2), top, Inches(1), Inches(0.4), pct, font_size=14, color=GRAY_400)

    # Milestones
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.5),
                 "18-Month Milestones to Series A:", font_size=16, bold=True, color=BRAND_BLUE)

    milestone_lines = [
        ("  ✔  15 paying customers ($570K ARR)", False, GREEN_400),
        ("  ✔  3 case studies with measurable ROI", False, GREEN_400),
        ("  ✔  Healthcare vertical dominance (5+ hospitals)", False, GREEN_400),
        ("  ✔  95% gross retention rate", False, GREEN_400),
        ("  ✔  Self-serve signup live (PLG engine running)", False, GREEN_400),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(5.5), Inches(8), Inches(2), milestone_lines, font_size=14)

    # ===================== SLIDE 13: Closing =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_rounded_rect(slide, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5), BRAND_BLUE, "SG", font_size=36)

    add_text_box(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(1),
                 "ShiftGuard", font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.2),
                 "Schedule with confidence.\nNever pay another compliance fine.",
                 font_size=22, color=GRAY_300, alignment=PP_ALIGN.CENTER)

    # Key line
    lines = [
        ("", False, None),
        ("Every schedule published without ShiftGuard is a gamble.", True, WHITE),
        ("We turn compliance from a liability into a competitive advantage.", False, BRAND_BLUE),
    ]
    add_multiline_text(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(1.5), lines, font_size=16)

    add_text_box(slide, Inches(3), Inches(6.6), Inches(7.3), Inches(0.5),
                 "hello@shiftguard.ai  |  frontend-lyart-three-52.vercel.app/demo",
                 font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "docs", "ShiftGuard_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"Deck saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_deck()
