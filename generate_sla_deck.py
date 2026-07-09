"""
ShiftGuard SLA & Support Model Presentation Generator
Generates a branded PowerPoint covering service tiers, uptime, support, security, and onboarding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BRAND_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
BRAND_DARK = RGBColor(0x03, 0x07, 0x12)
BRAND_NAVY = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GREEN_400 = RGBColor(0x4A, 0xDE, 0x80)
RED_400 = RGBColor(0xF8, 0x71, 0x71)
ORANGE_400 = RGBColor(0xFB, 0x92, 0x3C)
PURPLE_400 = RGBColor(0xC0, 0x84, 0xFC)
CYAN_400 = RGBColor(0x22, 0xD3, 0xEE)


def set_slide_bg(slide, color=BRAND_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    return box


def add_lines(slide, left, top, width, height, lines, size=14, line_spacing=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(size * line_spacing * 0.5)
    return box


def add_card(slide, left, top, width, height, fill=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def generate_sla_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================
    # SLIDE 1: Title
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "ShiftGuard", size=48, bold=True, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Service Level Agreement & Support Model", size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
             "What happens after you say yes.", size=18, color=GRAY_400, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "Confidential | July 2026", size=12, color=GRAY_500, align=PP_ALIGN.CENTER)

    # =========================================================
    # SLIDE 2: Service Tiers Overview
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Service Tiers", size=32, bold=True)

    tiers = [
        ("Starter", "$15", "25+", "99.5%", "Business hours", "Self-serve", GRAY_400),
        ("Professional", "$20", "50+", "99.9%", "Extended (7a-9p)", "Guided", BRAND_BLUE),
        ("Enterprise", "$25", "200+", "99.95%", "24/7", "White-glove", GREEN_400),
    ]
    for i, (name, price, headcount, uptime, support, onboard, accent) in enumerate(tiers):
        x = Inches(0.7 + i * 4.1)
        add_card(slide, x, Inches(1.4), Inches(3.8), Inches(5.5))
        add_text(slide, x + Inches(0.3), Inches(1.6), Inches(3.2), Inches(0.5),
                 name, size=22, bold=True, color=accent)
        add_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.5),
                 f"{price} PEPM", size=28, bold=True)
        add_lines(slide, x + Inches(0.3), Inches(3.0), Inches(3.4), Inches(3.5), [
            (f"Min headcount: {headcount}", False, GRAY_300),
            (f"Uptime SLA: {uptime}", False, GRAY_300),
            (f"Support: {support}", False, GRAY_300),
            (f"Onboarding: {onboard}", False, GRAY_300),
            (f"Law updates: {'7 days' if i==0 else '48 hours' if i==1 else '24 hours'}", False, GRAY_300),
            (f"Critical response: {'4h' if i==0 else '1h' if i==1 else '15 min'}", False, GRAY_300),
        ], size=13)

    # =========================================================
    # SLIDE 3: Uptime & Credits
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Uptime Guarantee & Service Credits", size=32, bold=True)

    add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
    add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "SLA Credit Schedule", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4), [
        ("99.9% - 99.95%    No credit (Enterprise meets SLA)", False, GREEN_400),
        ("99.5% - 99.9%      10% credit", False, GRAY_300),
        ("99.0% - 99.5%      25% credit", False, ORANGE_400),
        ("95.0% - 99.0%      50% credit", False, RED_400),
        ("Below 95.0%          100% credit (full month free)", False, RED_400),
        ("", False, GRAY_500),
        ("Max credit per month: 100% of fee", False, GRAY_400),
        ("Applied to following month's invoice", False, GRAY_400),
    ], size=14)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5))
    add_text(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.5),
             "Monitoring & Transparency", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(7.1), Inches(2.3), Inches(5.5), Inches(4), [
        ("Public status page: status.shiftguard.ai", False, GRAY_300),
        ("Synthetic monitoring every 60 seconds", False, GRAY_300),
        ("Incident alerts within 5 minutes of detection", False, GRAY_300),
        ("Post-incident reports within 72 hours", False, GRAY_300),
        ("", False, GRAY_500),
        ("Scheduled maintenance:", True, GRAY_300),
        ("Sundays 2-6 AM ET (72h advance notice)", False, GRAY_400),
        ("Max 4 hours/month planned downtime", False, GRAY_400),
    ], size=14)

    # =========================================================
    # SLIDE 4: Support Model
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Support Model", size=32, bold=True)

    add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(2.5))
    add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Response Times", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.5), [
        ("Critical (P1):   4h / 1h / 15 min", False, RED_400),
        ("High (P2):        8h / 4h / 1h", False, ORANGE_400),
        ("Normal (P3):    2 days / 1 day / 4h", False, GRAY_300),
    ], size=13)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(2.5))
    add_text(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Channels", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(1.5), [
        ("In-app AI (Otto):  All tiers", False, GRAY_300),
        ("Email:                    All tiers", False, GRAY_300),
        ("Phone:                  Pro + Enterprise", False, GRAY_300),
        ("Slack Connect:      Enterprise only", False, CYAN_400),
    ], size=13)

    add_card(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8))
    add_text(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
             "Severity Definitions", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2), [
        ("P1 Critical: Platform unavailable, incorrect compliance results, data loss, security breach", False, RED_400),
        ("P2 High: Major feature degraded, one jurisdiction not processing, >5s response times", False, ORANGE_400),
        ("P3 Normal: UI bugs, report formatting, minor performance, feature requests", False, GRAY_300),
        ("P4 Low: Cosmetic issues, enhancement requests, training questions", False, GRAY_500),
    ], size=13)

    # =========================================================
    # SLIDE 5: Compliance Rule Maintenance
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Compliance Rule Maintenance", size=32, bold=True)

    add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             "When labor laws change, your system updates automatically.", size=16, color=GRAY_400)

    add_card(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0))
    add_text(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.4),
             "Law Update SLA", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.2), [
        ("Enterprise:  24 hours from law change to deployed update", False, GREEN_400),
        ("Professional:  48 hours", False, BRAND_BLUE),
        ("Starter:  7 calendar days", False, GRAY_300),
    ], size=14)

    add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7))
    add_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
             "Update Process", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2), [
        ("1. Detection — AI scans federal/state registers, DOL bulletins daily", False, GRAY_300),
        ("2. Analysis — Compliance team assesses impact & affected jurisdictions", False, GRAY_300),
        ("3. Implementation — Rule engine updated, version-controlled, tested", False, GRAY_300),
        ("4. Notification — Affected customers notified (email + in-app banner)", False, GRAY_300),
        ("5. Audit trail — Previous rule version preserved for legal defense", False, GRAY_300),
    ], size=13)

    # =========================================================
    # SLIDE 6: Security & Data
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Security & Data Handling", size=32, bold=True)

    add_card(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3))
    add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "Infrastructure", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.2), [
        ("Hosted on AWS (us-east-1 + us-west-2 DR)", False, GRAY_300),
        ("Encryption at rest: AES-256", False, GRAY_300),
        ("Encryption in transit: TLS 1.3", False, GRAY_300),
        ("Secrets rotated every 90 days", False, GRAY_300),
        ("SOC 2 Type II (target Q2 2027)", False, GRAY_400),
    ], size=13)

    add_card(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3))
    add_text(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.4),
             "HIPAA (Healthcare Tier)", size=18, bold=True, color=GREEN_400)
    add_lines(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.2), [
        ("BAA available for Enterprise healthcare", False, GRAY_300),
        ("PHI isolation: dedicated database schema", False, GRAY_300),
        ("Access logging: every read/write audited", False, GRAY_300),
        ("Role-based access (minimum necessary)", False, GRAY_300),
        ("Annual third-party security audit", False, GRAY_300),
    ], size=13)

    add_card(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5))
    add_text(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
             "Data Ownership & Portability", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
        ("You own 100% of your data — always exportable (CSV, JSON, API)", False, GRAY_300),
        ("On termination: 30-day retrieval window, then permanent deletion + certificate", False, GRAY_300),
        ("No cross-customer data training without explicit written consent", False, GRAY_300),
        ("Breach notification: within 1 hour of confirmed incident", False, RED_400),
    ], size=13)

    # =========================================================
    # SLIDE 7: Onboarding Timeline
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Onboarding — Time to Value", size=32, bold=True)

    # Starter
    add_card(slide, Inches(0.5), Inches(1.4), Inches(3.9), Inches(5.5))
    add_text(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(0.4),
             "Starter: Same Day", size=18, bold=True, color=GRAY_400)
    add_lines(slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(4.5), [
        ("1. Sign up, select state", False, GRAY_300),
        ("2. Upload schedule (CSV/Excel)", False, GRAY_300),
        ("3. First report in < 60 seconds", False, GREEN_400),
        ("4. In-app walkthrough", False, GRAY_300),
        ("5. Otto AI for config questions", False, GRAY_300),
    ], size=13)

    # Professional
    add_card(slide, Inches(4.7), Inches(1.4), Inches(3.9), Inches(5.5))
    add_text(slide, Inches(5.0), Inches(1.6), Inches(3.5), Inches(0.4),
             "Professional: 1-2 Weeks", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(5.0), Inches(2.2), Inches(3.5), Inches(4.5), [
        ("1. Kickoff call (30 min)", False, GRAY_300),
        ("2. Config session (60 min)", False, GRAY_300),
        ("3. Go-live week: daily check-ins", False, GRAY_300),
        ("4. 30-day health check call", False, GRAY_300),
        ("", False, GRAY_500),
        ("First value: Day 1", True, GREEN_400),
    ], size=13)

    # Enterprise
    add_card(slide, Inches(8.9), Inches(1.4), Inches(3.9), Inches(5.5))
    add_text(slide, Inches(9.2), Inches(1.6), Inches(3.5), Inches(0.4),
             "Enterprise: 4-6 Weeks", size=18, bold=True, color=GREEN_400)
    add_lines(slide, Inches(9.2), Inches(2.2), Inches(3.5), Inches(4.5), [
        ("1. Discovery (1 week)", False, GRAY_300),
        ("2. Configuration (1-2 weeks)", False, GRAY_300),
        ("3. Parallel run (2 weeks)", False, GRAY_300),
        ("4. Go-live (1 day)", False, GREEN_400),
        ("5. Stabilization (4 weeks)", False, GRAY_300),
        ("6. Quarterly business reviews", False, GRAY_300),
        ("", False, GRAY_500),
        ("Includes CBA digitization", True, CYAN_400),
        ("UKG/Kronos integration", True, CYAN_400),
    ], size=13)

    # =========================================================
    # SLIDE 8: Disaster Recovery
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Disaster Recovery & Business Continuity", size=32, bold=True)

    metrics = [
        ("RPO", "1 hour", "Maximum data loss window"),
        ("RTO", "4h / 1h", "Starter-Pro / Enterprise"),
        ("Backups", "Hourly", "Incremental + daily full"),
        ("DR Region", "us-west-2", "Automatic failover"),
    ]
    for i, (label, value, desc) in enumerate(metrics):
        x = Inches(0.5 + i * 3.2)
        add_card(slide, x, Inches(1.4), Inches(2.9), Inches(2.5))
        add_text(slide, x + Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.4),
                 label, size=14, color=GRAY_400)
        add_text(slide, x + Inches(0.2), Inches(2.1), Inches(2.5), Inches(0.5),
                 value, size=26, bold=True, color=BRAND_BLUE)
        add_text(slide, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.4),
                 desc, size=12, color=GRAY_500)

    add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7))
    add_text(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
             "Release & Maintenance Schedule", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2), [
        ("Bug fixes:  Deployed as needed (zero-downtime rolling deploys)", False, GRAY_300),
        ("Features:  Bi-weekly (every other Tuesday)", False, GRAY_300),
        ("Major versions:  Quarterly (90-day migration window for breaking changes)", False, GRAY_300),
        ("API versioning:  Previous version supported 12 months after successor launches", False, GRAY_300),
    ], size=13)

    # =========================================================
    # SLIDE 9: Contract Summary
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             "Contract Terms at a Glance", size=32, bold=True)

    add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
    add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
             "Billing & Commitment", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("Monthly or Annual (10% discount)", False, GRAY_300),
        ("NET 30 payment terms", False, GRAY_300),
        ("No setup fees (Starter/Professional)", False, GREEN_400),
        ("Monthly: 30-day written notice to cancel", False, GRAY_300),
        ("Annual: auto-renews unless 60-day notice", False, GRAY_300),
        ("", False, GRAY_500),
        ("For-cause termination:", True, GRAY_300),
        ("Immediate if breach uncured after 30 days", False, GRAY_400),
        ("3 consecutive months SLA cap exceeded", False, GRAY_400),
    ], size=13)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5))
    add_text(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.4),
             "Data on Termination", size=18, bold=True, color=BRAND_BLUE)
    add_lines(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.5), [
        ("30-day data retrieval window", False, GRAY_300),
        ("Export: CSV, JSON, PDF, or via API", False, GRAY_300),
        ("Permanent deletion after 30 days", False, GRAY_300),
        ("Certificate of destruction provided", False, GRAY_300),
        ("", False, GRAY_500),
        ("No lock-in. Your data is yours.", True, GREEN_400),
    ], size=14)

    # =========================================================
    # SLIDE 10: Next Steps
    # =========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Ready to Get Started?", size=36, bold=True, align=PP_ALIGN.CENTER)

    add_lines(slide, Inches(3), Inches(3), Inches(7), Inches(3), [
        ("1.  Pick your tier (Starter / Professional / Enterprise)", False, GRAY_300),
        ("2.  We'll send the agreement for review", False, GRAY_300),
        ("3.  Sign + onboard (same day for Starter)", False, GRAY_300),
        ("4.  First compliance report in under 60 seconds", False, GREEN_400),
    ], size=16)

    add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "sales@shiftguard.ai  |  shiftguard.ai", size=16, color=BRAND_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
             "Schedule with confidence. Never pay another compliance fine.",
             size=14, color=GRAY_500, align=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "docs", "ShiftGuard_SLA_Support.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_sla_deck()
