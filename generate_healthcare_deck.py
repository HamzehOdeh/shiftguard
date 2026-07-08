"""
ShiftGuard for Healthcare — Sales Pitch Deck
Tailored for hospital administrators, CMOs, CNOs, Program Directors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BRAND_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
BRAND_DARK = RGBColor(0x03, 0x07, 0x12)
BRAND_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
RED_400 = RGBColor(0xF8, 0x71, 0x71)
GREEN_400 = RGBColor(0x4A, 0xDE, 0x80)
ORANGE_400 = RGBColor(0xFB, 0x92, 0x3C)
PURPLE_400 = RGBColor(0xA7, 0x8B, 0xFA)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, text_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color or color
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_stat_card(slide, left, top, number, label, color=BRAND_BLUE):
    width = Inches(2.2)
    height = Inches(1.4)
    shape = add_rounded_rect(slide, left, top, width, height, BRAND_NAVY)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_400
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER


def create_healthcare_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== SLIDE 1: Title =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    logo = add_rounded_rect(slide, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5), BRAND_BLUE, "SG", font_size=36)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.2),
                 "ShiftGuard for Healthcare", font_size=44, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.8),
                 "AI-Powered Clinical Scheduling Compliance", font_size=22, color=GRAY_300, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.1), Inches(9.3), Inches(0.8),
                 "Stop scheduling residents into ACGME violations.\nStop losing nurses to unfair schedules.",
                 font_size=16, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(4), Inches(6.5), Inches(5.3), Inches(0.5),
                 "July 2026 | For Hospital Leadership", font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    # ===================== SLIDE 2: The Healthcare Problem =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "The Problem: Clinical Scheduling Is a Compliance Minefield", font_size=28, bold=True)

    # Stat cards - healthcare specific
    add_stat_card(slide, Inches(0.5), Inches(1.5), "6 months", "Avg ACGME\nprobation length", RED_400)
    add_stat_card(slide, Inches(3.0), Inches(1.5), "$7,500", "Per ACGME\nviolation/week", ORANGE_400)
    add_stat_card(slide, Inches(5.5), Inches(1.5), "43%", "Nurse turnover\nfrom unfair schedules", RED_400)
    add_stat_card(slide, Inches(8.0), Inches(1.5), "$56K", "Cost to replace\none RN", ORANGE_400)
    add_stat_card(slide, Inches(10.5), Inches(1.5), "4-8 hrs", "Chief Resident time/wk\non schedule compliance", GRAY_300)

    lines = [
        ("What keeps hospital leaders up at night:", True, WHITE),
        ("", False, None),
        ("  Program Directors: \"Did any resident exceed 80 hours this week?\"", False, RED_400),
        ("  — ACGME violations can shut down your residency program", False, GRAY_400),
        ("", False, None),
        ("  CNOs: \"Are we meeting nurse-to-patient ratios on every unit?\"", False, RED_400),
        ("  — CMS citations, Joint Commission findings, patient safety events", False, GRAY_400),
        ("", False, None),
        ("  CMOs: \"Is our call schedule fair — or is it a burnout machine?\"", False, RED_400),
        ("  — Unfair schedules drive 43% of nursing turnover ($56K per nurse lost)", False, GRAY_400),
        ("", False, None),
        ("  HR/Compliance: \"Can we prove we checked before we published?\"", False, RED_400),
        ("  — DOL investigations, class-actions, no audit trail of compliance checks", False, GRAY_400),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(3.2), Inches(12), Inches(4), lines, font_size=13)

    # ===================== SLIDE 3: Real Scenarios =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "This Happens Every Week in Your Hospital", font_size=28, bold=True)

    # Scenario cards
    scenarios = [
        ("ACGME Violation — Caught Too Late",
         "Dr. Patel worked 84 hours this week. Nobody noticed until the quarterly ACGME report. The program is now on probation for 6 months. Three applicants withdrew their match preferences.",
         "ShiftGuard flags it Monday morning — BEFORE the schedule is published.",
         RED_400),
        ("Nurse Clopening — Fatigue Incident",
         "RN Sarah Chen closed at 11PM Saturday, opened at 6AM Sunday (5hr rest). She made a medication error at 4AM. The family filed a complaint. CMS is investigating.",
         "ShiftGuard blocks the schedule: \"REST-002: Less than 8h rest. Fix: Move Sunday start to 10AM.\"",
         ORANGE_400),
        ("Unfair Call Distribution — Grievance Filed",
         "The nurses' union filed a grievance: 4 nurses worked every holiday in 2025 while 6 never worked one. The chief always gave herself preferred schedules. Arbitration cost $180K.",
         "ShiftGuard's fairness engine distributes automatically. Every decision explained in writing.",
         PURPLE_400),
    ]

    for i, (title, problem, solution, color) in enumerate(scenarios):
        top = Inches(1.4 + i * 2.0)
        add_rounded_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.8), BRAND_NAVY)
        add_rounded_rect(slide, Inches(0.8), top, Inches(0.15), Inches(1.8), color)
        add_text_box(slide, Inches(1.3), top + Inches(0.1), Inches(10), Inches(0.4), title, font_size=14, bold=True, color=color)
        add_text_box(slide, Inches(1.3), top + Inches(0.5), Inches(10.5), Inches(0.5), problem, font_size=11, color=GRAY_400)
        sol_text = solution.split("ShiftGuard ")[1] if "ShiftGuard " in solution else solution
        add_text_box(slide, Inches(1.3), top + Inches(1.2), Inches(10.5), Inches(0.5), sol_text, font_size=11, color=GREEN_400, bold=True)

    # ===================== SLIDE 4: How It Works =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "How ShiftGuard Works for Hospitals", font_size=28, bold=True)

    # Role-based views
    roles = [
        ("Program Director", [
            "• See ACGME compliance status across ALL residents at a glance",
            "• 80h/week 4-week rolling average — tracked automatically",
            "• Swap requests pre-checked: \"Will this swap violate 24+4?\"",
            "• Generate fair year schedule (AI distributes nights/weekends equally)",
            "• One-click ACGME report for site visits",
        ], RGBColor(0x06, 0x94, 0xA2)),
        ("Nurse Manager / CNO", [
            "• Staffing ratios checked per unit (ICU 1:2, Med-Surg 1:5, ED 1:4)",
            "• Clopening detection with auto-fix suggestions",
            "• Fairness dashboard: OT, nights, weekends, holidays per nurse",
            "• Self-healing: callout at 5AM → auto-finds qualified, rested coverage",
            "• Credential expiry alerts (BLS, ACLS, RN license, NRP)",
        ], BRAND_BLUE),
        ("Admin / HR / Compliance", [
            "• Full audit trail: every schedule checked, timestamped, exportable",
            "• Bias audit (EEOC four-fifths rule, NYC Local Law 144)",
            "• FMLA auto-triggers on 3+ day absences (within 5-day federal deadline)",
            "• DOL investigation defense package (one-click)",
            "• ROI dashboard: penalties avoided, retention savings, time saved",
        ], PURPLE_400),
    ]

    for i, (role, bullets, color) in enumerate(roles):
        left = Inches(0.5 + i * 4.2)
        add_rounded_rect(slide, left, Inches(1.5), Inches(4.0), Inches(5.3), BRAND_NAVY)
        add_text_box(slide, left + Inches(0.3), Inches(1.6), Inches(3.5), Inches(0.5), role, font_size=15, bold=True, color=color)
        bullet_lines = [(b, False, GRAY_300) for b in bullets]
        add_multiline_text(slide, left + Inches(0.3), Inches(2.2), Inches(3.6), Inches(4.5), bullet_lines, font_size=11, spacing=1.8)

    # ===================== SLIDE 5: ACGME Deep Dive =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "ACGME Compliance — Automated, Not Manual", font_size=28, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
                 "Every Common Program Requirement checked in real-time for every resident:", font_size=15, color=GRAY_300)

    # ACGME rules
    rules = [
        ("80-Hour Weekly Limit", "4-week rolling average, including moonlighting", "✔ Auto-tracked"),
        ("24+4 Hour Shift Cap", "Max 24h continuous + 4h handoff/education", "✔ Flagged at hour 20"),
        ("1 Day Off in 7", "Averaged over 4 weeks", "✔ Auto-calculated"),
        ("10-Hour Rest Between Shifts", "After 24h duty or night float", "✔ Blocks violations"),
        ("14-Hour Shift Limit (PGY-1)", "Interns cannot exceed 14h continuous", "✔ PGY-aware"),
        ("Night Float: Max 6 Consecutive", "Cannot exceed 6 consecutive night shifts", "✔ Auto-counts"),
        ("Every 3rd Night Call Max", "In-house call frequency limit", "✔ Pattern detection"),
        ("4h/Week Education Protected", "Didactic time cannot be scheduled over", "✔ Blocks scheduling"),
        ("Moonlighting Counts Toward 80h", "External work included in cap", "✔ Self-reported + tracked"),
    ]

    headers = ["ACGME Rule", "What It Means", "ShiftGuard"]
    col_widths = [Inches(3.5), Inches(5.5), Inches(2.5)]
    start_left = Inches(0.8)
    start_top = Inches(1.9)

    # Header
    x = start_left
    for i, h in enumerate(headers):
        add_text_box(slide, x, start_top, col_widths[i], Inches(0.4), h, font_size=11, bold=True, color=BRAND_BLUE)
        x += col_widths[i]

    for r_idx, (rule, meaning, status) in enumerate(rules):
        y = start_top + Inches(0.45) + Inches(0.52) * r_idx
        x = start_left
        add_text_box(slide, x, y, col_widths[0], Inches(0.4), rule, font_size=11, bold=True, color=WHITE)
        add_text_box(slide, x + col_widths[0], y, col_widths[1], Inches(0.4), meaning, font_size=10, color=GRAY_400)
        add_text_box(slide, x + col_widths[0] + col_widths[1], y, col_widths[2], Inches(0.4), status, font_size=10, color=GREEN_400)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.4),
                 "These rules are LOCKED — cannot be overridden. ShiftGuard enforces them as law, not preference.",
                 font_size=12, bold=True, color=ORANGE_400)

    # ===================== SLIDE 6: AI Assistant Otto =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Meet Otto — Your AI Scheduling Assistant", font_size=28, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
                 "Ask anything about your schedule in plain English. Get instant, accurate answers.", font_size=15, color=GRAY_300)

    # Chat mockup
    add_rounded_rect(slide, Inches(0.8), Inches(1.9), Inches(7), Inches(5.0), BRAND_NAVY)

    chat_lines = [
        ("You: \"Can Dr. Patel cover tonight's shift?\"", False, GRAY_300),
        ("", False, None),
        ("Otto: No. Dr. Patel is at 76 hours this week (4-week avg: 78h).", False, WHITE),
        ("Adding tonight's 12h shift would put him at 88h — exceeding the", False, WHITE),
        ("ACGME 80h limit by 8 hours.", False, WHITE),
        ("", False, None),
        ("Recommended alternative: Dr. Kim (62h this week, lowest night", False, GREEN_400),
        ("count this rotation, ACGME-safe for 12h coverage).", False, GREEN_400),
        ("", False, None),
        ("You: \"Generate a fair call schedule for July\"", False, GRAY_300),
        ("", False, None),
        ("Otto: Done. 31 nights distributed across 5 residents:", False, WHITE),
        ("  • Patel: 6 nights (was 7 last month — rebalanced)", False, WHITE),
        ("  • Kim: 7 nights | Reeves: 6 | Santos: 6 | Park: 6", False, WHITE),
        ("  Night deviation: 0.4 (excellent fairness).", False, GREEN_400),
        ("  No ACGME violations. No 24+4 conflicts.", False, GREEN_400),
        ("  [Publish]  [Adjust]  [Show Fairness Scorecard]", False, BRAND_BLUE),
    ]
    add_multiline_text(slide, Inches(1.1), Inches(2.1), Inches(6.5), Inches(4.5), chat_lines, font_size=11, spacing=1.1)

    # Side benefits
    benefit_lines = [
        ("Otto knows:", True, BRAND_BLUE),
        ("", False, None),
        ("• Every resident's hours (live)", False, GRAY_300),
        ("• ACGME rules per PGY level", False, GRAY_300),
        ("• Who's on call, on leave, available", False, GRAY_300),
        ("• Historical fairness scores", False, GRAY_300),
        ("• Credential status", False, GRAY_300),
        ("• CBA / union constraints", False, GRAY_300),
        ("", False, None),
        ("No clarifying questions.", True, ORANGE_400),
        ("Answers directly with data.", True, ORANGE_400),
        ("", False, None),
        ("Available 24/7.", False, GRAY_400),
        ("Cheaper than a phone call", False, GRAY_400),
        ("to the chief at 2AM.", False, GRAY_400),
    ]
    add_multiline_text(slide, Inches(8.3), Inches(1.9), Inches(4.5), Inches(5), benefit_lines, font_size=13, spacing=1.3)

    # ===================== SLIDE 7: Jeopardy System =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Self-Healing Schedules: Jeopardy + Predictive Callouts", font_size=28, bold=True)

    # Left: jeopardy system
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "Jeopardy (Backup) System", font_size=16, bold=True, color=RGBColor(0x06, 0x94, 0xA2))

    jeopardy_lines = [
        ("How it works today (manual):", True, GRAY_400),
        ("  Resident calls out at 5AM → Chief gets woken up →", False, GRAY_400),
        ("  Calls 4 people → finds someone → 2 hours of chaos", False, GRAY_400),
        ("", False, None),
        ("How it works with ShiftGuard:", True, WHITE),
        ("", False, None),
        ("  1. Pre-assigned backup for every shift (fairness-ranked)", False, GREEN_400),
        ("  2. ACGME pre-checked (won't assign if it violates 80h)", False, GREEN_400),
        ("  3. Callout confirmed → backup auto-activated", False, GREEN_400),
        ("  4. Chief wakes up to: \"Covered. Dr. Kim accepted 5:12AM\"", False, GREEN_400),
        ("", False, None),
        ("Pattern learning:", True, BRAND_BLUE),
        ("  • Learns who calls out on which days", False, GRAY_300),
        ("  • Post-night-float, post-call patterns detected", False, GRAY_300),
        ("  • 24-48h lookahead: \"78% chance callout Friday\"", False, GRAY_300),
        ("  • Pre-positions standby BEFORE the callout happens", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(1.1), Inches(2.1), Inches(5.5), Inches(4.5), jeopardy_lines, font_size=12, spacing=1.2)

    # Right: nursing coverage
    add_rounded_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(7.5), Inches(1.6), Inches(5), Inches(0.5),
                 "Nursing Self-Healing", font_size=16, bold=True, color=BRAND_BLUE)

    nursing_lines = [
        ("When a gap appears:", True, WHITE),
        ("", False, None),
        ("  Step 1: Check pre-declared volunteers first", False, GRAY_300),
        ("  Step 2: Offer to lowest-OT qualified nurses", False, GRAY_300),
        ("  Step 3: Broadcast to all eligible (Uber model)", False, GRAY_300),
        ("  Step 4: Mandatory assignment (fairness-ranked)", False, GRAY_300),
        ("  Step 5: Agency escalation (last resort)", False, GRAY_300),
        ("", False, None),
        ("Every step checks:", True, BRAND_BLUE),
        ("  • Rest period compliance", False, GREEN_400),
        ("  • OT budget impact", False, GREEN_400),
        ("  • Credential validity", False, GREEN_400),
        ("  • Unit competency (ICU-trained? ED-trained?)", False, GREEN_400),
        ("  • Fairness score (who's been asked most?)", False, GREEN_400),
        ("", False, None),
        ("Result: Gaps filled in minutes, not hours.", True, GREEN_400),
        ("Manager time saved: 5-10 hours/week.", False, GRAY_400),
    ]
    add_multiline_text(slide, Inches(7.5), Inches(2.1), Inches(5), Inches(4.5), nursing_lines, font_size=12, spacing=1.2)

    # ===================== SLIDE 8: Fairness Engine =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Fairness Engine — Your Retention Weapon", font_size=28, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
                 "43% of nurse turnover is driven by perceived schedule unfairness. Each lost nurse costs $56K.",
                 font_size=14, color=ORANGE_400)

    # What fairness tracks
    fairness_items = [
        ("Night Shifts", "Equal distribution across rotation. Year-over-year tracking."),
        ("Weekends", "Nobody works 3 weekends in a row while others get every weekend off."),
        ("Holidays", "Priority-based auction. If you got Christmas, you're lower priority for New Year's."),
        ("Overtime", "OT hours distributed equitably. No one person always gets asked."),
        ("Undesirable Shifts", "Friday nights, holiday eves, short-notice coverage — tracked and balanced."),
        ("Request Burden", "Who gets called for coverage most? System rebalances automatically."),
    ]

    for i, (title, desc) in enumerate(fairness_items):
        col = i % 2
        row = i // 2
        left = Inches(0.8 + col * 6.3)
        top = Inches(1.9 + row * 1.6)
        add_rounded_rect(slide, left, top, Inches(5.9), Inches(1.3), BRAND_NAVY)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(5), Inches(0.4), title, font_size=14, bold=True, color=GREEN_400)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.55), Inches(5.3), Inches(0.7), desc, font_size=12, color=GRAY_300)

    # Bottom: key differentiator
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
                 "Every decision explained in plain English: \"Assigned to Sarah because she has the lowest night count and is ACGME-safe.\"",
                 font_size=13, color=BRAND_BLUE, bold=True)

    # ===================== SLIDE 9: vs Competitors =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Why Not Just Use QGenda / Amion / UKG?", font_size=28, bold=True)

    headers = ["", "ShiftGuard", "QGenda", "Amion", "UKG"]
    rows = [
        ["Real-time ACGME tracking", "✔", "✖ Retrospective", "✖ Retrospective", "✖"],
        ["Pre-publish violation check", "✔", "✖", "✖", "✖"],
        ["AI fix suggestions", "✔", "✖", "✖", "✖"],
        ["Penalty $ exposure", "✔", "✖", "✖", "✖"],
        ["Fairness engine", "✔", "✖", "✖", "✖"],
        ["Jeopardy auto-activation", "✔", "Partial", "✖", "✖"],
        ["Plain English explanations", "✔", "✖", "✖", "✖"],
        ["Nurse + Resident in one tool", "✔", "Resident only", "Resident only", "Nurses only"],
        ["Setup time", "1 hour", "2-4 weeks", "1-2 weeks", "6-12 months"],
        ["Price", "$20 PEPM", "$30-50 PEPM", "$15/provider", "$50-100+ PEPM"],
    ]

    col_widths = [Inches(3.2), Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.2)]
    start_top = Inches(1.4)
    row_height = Inches(0.52)
    start_left = Inches(0.5)

    x = start_left
    for i, h in enumerate(headers):
        color = BRAND_BLUE if i == 1 else GRAY_400
        add_text_box(slide, x, start_top, col_widths[i], Inches(0.4), h, font_size=11, bold=True, color=color)
        x += col_widths[i]

    for r_idx, row in enumerate(rows):
        y = start_top + Inches(0.5) + row_height * r_idx
        x = start_left
        for c_idx, cell in enumerate(row):
            if cell == "✔":
                color = GREEN_400
            elif cell == "✖" or cell.startswith("✖"):
                color = RED_400
            elif c_idx == 1:
                color = GREEN_400
            elif c_idx == 0:
                color = WHITE
            else:
                color = GRAY_400
            add_text_box(slide, x, y, col_widths[c_idx], row_height, cell, font_size=11, color=color)
            x += col_widths[c_idx]

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
                 "Bottom line: QGenda/Amion tell you AFTER. We prevent BEFORE. And we cover nursing + residents in one platform.",
                 font_size=13, color=BRAND_BLUE, bold=True)

    # ===================== SLIDE 10: ROI =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "ROI: This Pays for Itself in Week 1", font_size=28, bold=True)

    # ROI calculation
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "For a 200-Bed Hospital (300 clinical staff)", font_size=14, bold=True, color=BRAND_BLUE)

    roi_lines = [
        ("Annual Cost of ShiftGuard:", True, WHITE),
        ("  300 staff × $20/mo × 12 = $72,000/year", False, GRAY_300),
        ("", False, None),
        ("Annual Savings:", True, GREEN_400),
        ("", False, None),
        ("  ACGME violations avoided (2/month × $7,500)     $180,000", False, GREEN_400),
        ("  Nurse turnover reduction (4 fewer/year × $56K)  $224,000", False, GREEN_400),
        ("  Agency spend reduction (20% less travel nurses)  $340,000", False, GREEN_400),
        ("  Manager time saved (6h/week × 10 managers)       $93,600", False, GREEN_400),
        ("  DOL/CMS penalty avoidance                        $50,000", False, GREEN_400),
        ("", False, None),
        ("  Total Annual Savings:                           $887,600", True, GREEN_400),
        ("", False, None),
        ("  ROI: 12.3x  |  Payback: 30 days", True, WHITE),
    ]
    add_multiline_text(slide, Inches(1.1), Inches(2.1), Inches(5.5), Inches(4.5), roi_lines, font_size=12, spacing=1.2)

    # Right side: qualitative
    add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.3), BRAND_NAVY)
    add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
                 "Beyond Dollars", font_size=14, bold=True, color=PURPLE_400)

    qual_lines = [
        ("Accreditation Protection", True, WHITE),
        ("  ACGME compliance 100% of the time", False, GRAY_300),
        ("  Audit-ready documentation always current", False, GRAY_300),
        ("", False, None),
        ("Patient Safety", True, WHITE),
        ("  No fatigued providers on duty", False, GRAY_300),
        ("  Proper staffing ratios enforced", False, GRAY_300),
        ("  Credential expiry never missed", False, GRAY_300),
        ("", False, None),
        ("Staff Satisfaction", True, WHITE),
        ("  Transparent, provably fair schedules", False, GRAY_300),
        ("  Self-service requests (auto-approved)", False, GRAY_300),
        ("  Mobile app for schedule access anywhere", False, GRAY_300),
        ("", False, None),
        ("Legal Protection", True, WHITE),
        ("  Every compliance check timestamped", False, GRAY_300),
        ("  Bias audit exportable for EEOC", False, GRAY_300),
        ("  Court-admissible audit trail", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(4.5), qual_lines, font_size=12, spacing=1.1)

    # ===================== SLIDE 11: Implementation =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Implementation: Live in 1 Week, Not 6 Months", font_size=28, bold=True)

    steps = [
        ("Day 1", "Setup", "Add your residents/nurses, define rotations,\nset department structure. 1-hour guided wizard.", BRAND_BLUE),
        ("Day 2-3", "Connect", "Link to UKG/Kronos (API) or upload schedule CSV.\nHistorical data imported for fairness baseline.", RGBColor(0x06, 0x94, 0xA2)),
        ("Day 4-5", "Validate", "Run compliance check on current schedule.\nSee violations that exist TODAY. Fix them.", PURPLE_400),
        ("Day 6-7", "Go Live", "Enable pre-publish checks. Workers get mobile app.\nOtto available for scheduling questions 24/7.", GREEN_400),
    ]

    for i, (day, title, desc, color) in enumerate(steps):
        top = Inches(1.5 + i * 1.4)
        add_rounded_rect(slide, Inches(0.8), top, Inches(1.5), Inches(1.1), color, day, font_size=14)
        add_text_box(slide, Inches(2.6), top + Inches(0.05), Inches(3), Inches(0.4), title, font_size=16, bold=True, color=color)
        add_text_box(slide, Inches(2.6), top + Inches(0.45), Inches(9), Inches(0.7), desc, font_size=12, color=GRAY_300)

    # What we DON'T require
    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.4),
                 "What we DON'T require:", font_size=14, bold=True, color=ORANGE_400)
    no_lines = [
        ("  ✖  No hardware installation", False, GRAY_300),
        ("  ✖  No IT project (cloud-based, nothing to install)", False, GRAY_300),
        ("  ✖  No replacing your existing systems (we sit on top)", False, GRAY_300),
        ("  ✖  No 6-month implementation consultants", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(6.6), Inches(8), Inches(1), no_lines, font_size=11, spacing=1.0)

    # ===================== SLIDE 12: Live Demo =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "See It Live — Right Now", font_size=28, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
                 "No slides needed. Open your browser:", font_size=16, color=GRAY_300)

    # URL cards
    urls = [
        ("Healthcare App (Full Product)", "shiftguard-giznml7xnt59j5aguqjerc.streamlit.app", "See ACGME compliance, residency scheduling,\nnursing management, Otto AI, bias audit, ROI dashboard"),
        ("Worker Mobile App", "frontend-lyart-three-52.vercel.app/worker", "What your nurses/residents see on their phone.\nSchedule, requests, balances, notifications."),
        ("Free Compliance Calculator", "frontend-lyart-three-52.vercel.app/calculator", "Enter your state + headcount → instant risk grade.\nShare with any colleague (no signup required)."),
    ]

    for i, (title, url, desc) in enumerate(urls):
        top = Inches(2.0 + i * 1.7)
        add_rounded_rect(slide, Inches(0.8), top, Inches(11.7), Inches(1.4), BRAND_NAVY)
        add_text_box(slide, Inches(1.3), top + Inches(0.1), Inches(5), Inches(0.4), title, font_size=15, bold=True, color=BRAND_BLUE)
        add_text_box(slide, Inches(1.3), top + Inches(0.5), Inches(5), Inches(0.4), url, font_size=13, color=GREEN_400)
        add_text_box(slide, Inches(7.0), top + Inches(0.2), Inches(5), Inches(1), desc, font_size=12, color=GRAY_400)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
                 "\"I'll show you — it takes 60 seconds. Pull up your phone.\"",
                 font_size=16, bold=True, color=WHITE)

    # ===================== SLIDE 13: Next Steps =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Next Steps", font_size=32, bold=True)

    steps_lines = [
        ("We're looking for 3 design partners — hospitals that want to:", True, WHITE),
        ("", False, None),
        ("  1.  Run ShiftGuard on their actual schedule data (free pilot)", False, GREEN_400),
        ("  2.  Measure real violation reduction over 90 days", False, GREEN_400),
        ("  3.  Become a published case study when results are proven", False, GREEN_400),
        ("", False, None),
        ("What you get as a design partner:", True, BRAND_BLUE),
        ("", False, None),
        ("  •  Free for 90 days (no commitment until ROI proven)", False, GRAY_300),
        ("  •  Direct access to the product team (features built to your needs)", False, GRAY_300),
        ("  •  Locked-in pricing (50% off standard rate, forever)", False, GRAY_300),
        ("  •  First to get new features (jeopardy AI, demand sensing, etc.)", False, GRAY_300),
        ("", False, None),
        ("Ideal partner:", True, WHITE),
        ("  • 100-500 clinical staff (residents + nurses)", False, GRAY_300),
        ("  • Residency program (ACGME compliance pain)", False, GRAY_300),
        ("  • Currently using UKG, QGenda, Amion, or spreadsheets", False, GRAY_300),
        ("  • Willing to share de-identified outcome metrics", False, GRAY_300),
    ]
    add_multiline_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(5.5), steps_lines, font_size=14, spacing=1.2)

    # ===================== SLIDE 14: Closing =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_DARK)

    add_rounded_rect(slide, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5), BRAND_BLUE, "SG", font_size=36)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1),
                 "ShiftGuard for Healthcare", font_size=40, bold=True, alignment=PP_ALIGN.CENTER)

    lines = [
        ("Your residents shouldn't be working illegal hours.", False, WHITE),
        ("Your nurses shouldn't quit over unfair schedules.", False, WHITE),
        ("Your program shouldn't be one violation away from probation.", False, WHITE),
        ("", False, None),
        ("We fix all three. Automatically. Before you publish.", True, BRAND_BLUE),
    ]
    add_multiline_text(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(2), lines, font_size=17, spacing=1.5)

    add_text_box(slide, Inches(3), Inches(6.5), Inches(7.3), Inches(0.5),
                 "hello@shiftguard.ai  |  Let's schedule a 15-minute walkthrough",
                 font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "docs", "ShiftGuard_Healthcare_Pitch.pptx")
    prs.save(output_path)
    print(f"Healthcare deck saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_healthcare_deck()
